# ai.py

import numpy as np
import cv2

# for file-processing LLM
import openai
from PyPDF2 import PdfReader
import fitz  # For PDF
from pptx import Presentation  # For PPTX
import docx  # For DOCX
import pandas as pd
import io
from PIL import Image
import base64

def analyze_candlestick_section(image):
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    thresh = cv2.adaptiveThreshold(blurred, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV, 11, 2)
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    candlestick_centers = []
    for cnt in contours:
        x, y, w, h = cv2.boundingRect(cnt)
        center_y = y + h // 2
        candlestick_centers.append((x + w // 2, center_y))
    candlestick_centers = sorted(candlestick_centers, key=lambda x: x[0])
    x_coords = [point[0] for point in candlestick_centers]
    y_coords = [point[1] for point in candlestick_centers]
    if len(x_coords) > 1:
        poly_fit = np.polyfit(x_coords, y_coords, 1)
        trend_slope = poly_fit[0]
        if trend_slope < 0:
            trend = "Upward trend"
        elif trend_slope > 0:
            trend = "Downward trend"
        else:
            trend = "Sideways trend"
    else:
        trend = "Not enough candlesticks"
    return trend

def analyze_graph_sections(image_path, num_sections=3):
    image = cv2.imread(image_path)
    image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    height, width = image.shape[:2]
    section_width = width // num_sections
    sections = [(i * section_width, (i + 1) * section_width) for i in range(num_sections)]
    section_trends = []
    cropped_images = []
    for section in sections:
        cropped_image = image[:, section[0]:section[1]]
        cropped_images.append(cropped_image)
        section_trend = analyze_candlestick_section(cropped_image)
        section_trends.append(section_trend)
    return section_trends, cropped_images

# for file processing LLM
# Set up OpenAI API key
openai.api_key = 'sk-proj-jxiVl2buV6WfIODqxThQT3BlbkFJn1bViSVnZDrzlNx6qrh2'

def read_document(file):
    text = ""
    file_content = file.read()
    filename = file.name

    if filename.endswith('.pdf'):
        with fitz.open(stream=file_content, filetype='pdf') as doc:
            for page in doc:
                text += page.get_text()
    elif filename.endswith('.docx'):
        doc = docx.Document(io.BytesIO(file_content))
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    elif filename.endswith('.txt'):
        text += file_content.decode('utf-8')
    elif filename.endswith('.pptx'):
        presentation = Presentation(io.BytesIO(file_content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif filename.endswith('.csv'):
        df = pd.read_csv(io.BytesIO(file_content))
        text += df.to_string()
    elif filename.endswith(('.xls', '.xlsx')):
        df = pd.read_excel(io.BytesIO(file_content))
        text += df.to_string()
    elif filename.endswith('.json'):
        data = pd.read_json(io.BytesIO(file_content))
        text += data.to_string()
    else:
        raise Exception("Unsupported file format.")
    
    return text

def split_text(text, max_tokens):
    words = text.split()
    chunks = []
    current_chunk = []
    current_length = 0

    for word in words:
        current_length += len(word) + 1  # +1 for space
        if current_length > max_tokens:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_length = len(word) + 1
        else:
            current_chunk.append(word)

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks

def ask_questionEd_about_document(document_chunks, question):
    answers = []
    for chunk in document_chunks:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an assistant who can answer questions about documents."},
                {"role": "user", "content": f"Document chunk: {chunk}\n\nQuestion: {question}"}
            ],
            max_tokens=1000
        )
        answers.append(response.choices[0].message['content'].strip())
    return " ".join(answers)

def encode_image(image):
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='PNG')
    return base64.b64encode(img_byte_arr.getvalue()).decode('utf-8')

def ask_question_about_image(question, image_url):
    response = openai.ChatCompletion.create(
        model='gpt-4-turbo',
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": question},
                    {"type": "image_url", "image_url": {"url": image_url}}
                ],
            }
        ],
        max_tokens=500,
    )

    answer = response.choices[0].message['content']
    return answer
